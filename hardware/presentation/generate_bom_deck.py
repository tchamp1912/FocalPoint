#!/usr/bin/env python3
"""Generate the editable Rev A BOM presentation."""

from pathlib import Path
import sys

sys.path.insert(0, "/private/tmp/focalpoint-pptx")
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "hardware" / "presentation" / "focalpoint-rev-a-bom.pptx"
PREVIEW = ROOT / "case" / "output" / "focalpoint-preview.png"

INK = RGBColor(25, 28, 36)
MUTED = RGBColor(96, 103, 115)
BLUE = RGBColor(57, 111, 237)
PALE = RGBColor(237, 243, 255)
GREEN = RGBColor(35, 145, 94)
AMBER = RGBColor(215, 137, 35)
WHITE = RGBColor(255, 255, 255)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def title(slide, text, subtitle=None):
    box = slide.shapes.add_textbox(Inches(.7), Inches(.45), Inches(11.9), Inches(.7))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.name = "Avenir Next"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INK
    if subtitle:
        sb = slide.shapes.add_textbox(Inches(.72), Inches(1.12), Inches(11.7), Inches(.45))
        sp = sb.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = "Avenir Next"
        sp.font.size = Pt(12)
        sp.font.color.rgb = MUTED


def footer(slide, number):
    line = slide.shapes.add_shape(1, Inches(.7), Inches(7.12), Inches(11.9), Inches(.012))
    line.fill.solid(); line.fill.fore_color.rgb = RGBColor(220, 224, 232); line.line.fill.background()
    b = slide.shapes.add_textbox(Inches(11.7), Inches(7.15), Inches(.8), Inches(.2))
    p = b.text_frame.paragraphs[0]; p.text = str(number); p.alignment = PP_ALIGN.RIGHT
    p.font.name = "Avenir Next"; p.font.size = Pt(9); p.font.color.rgb = MUTED


def bullets(slide, items, x, y, w, h, size=18, color=INK):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item; p.level = 0; p.font.name = "Avenir Next"; p.font.size = Pt(size)
        p.font.color.rgb = color; p.space_after = Pt(11)
    return box


def card(slide, x, y, w, h, heading, body, accent=BLUE):
    s = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = WHITE; s.line.color.rgb = RGBColor(216, 222, 232)
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(.07), Inches(h))
    bar.fill.solid(); bar.fill.fore_color.rgb = accent; bar.line.fill.background()
    hb = slide.shapes.add_textbox(Inches(x+.25), Inches(y+.18), Inches(w-.45), Inches(.35))
    hp = hb.text_frame.paragraphs[0]; hp.text = heading; hp.font.name = "Avenir Next"
    hp.font.size = Pt(15); hp.font.bold = True; hp.font.color.rgb = INK
    bb = slide.shapes.add_textbox(Inches(x+.25), Inches(y+.62), Inches(w-.45), Inches(h-.78))
    bb.text_frame.word_wrap = True
    bp = bb.text_frame.paragraphs[0]; bp.text = body; bp.font.name = "Avenir Next"
    bp.font.size = Pt(11); bp.font.color.rgb = MUTED


# 1 — title
s = prs.slides.add_slide(prs.slide_layouts[6])
shape = s.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
shape.fill.solid(); shape.fill.fore_color.rgb = INK; shape.line.fill.background()
tb = s.shapes.add_textbox(Inches(.85), Inches(1.1), Inches(7.3), Inches(1.5))
p = tb.text_frame.paragraphs[0]; p.text = "FocalPoint Rev A"; p.font.name = "Avenir Next"; p.font.size = Pt(44); p.font.bold = True; p.font.color.rgb = WHITE
p2 = tb.text_frame.add_paragraph(); p2.text = "Two-device bill of materials"; p2.font.name = "Avenir Next"; p2.font.size = Pt(25); p2.font.color.rgb = RGBColor(182, 204, 255)
bullets(s, ["16 dynamically mapped inputs", "Wireless + USB-C • 1,000 mAh • per-key RGB", "Prototype procurement decision — 28 July 2026"], .9, 3.2, 6.8, 2.1, 17, WHITE)
if PREVIEW.exists():
    s.shapes.add_picture(str(PREVIEW), Inches(8.2), Inches(1.0), width=Inches(4.4))

# 2 — product definition
s = prs.slides.add_slide(prs.slide_layouts[6]); title(s, "What we are building", "The BOM follows the current square enclosure and 16-input architecture")
card(s,.7,1.75,3.8,2.0,"Controls","13 tactile MX keys\n1 analog joystick + push\n1 rotary encoder + push\n1 capacitive touch region",BLUE)
card(s,4.75,1.75,3.8,2.0,"Industrial design","114 × 114 mm square shell\n6 mm corner radius\n4° key-plane slope\n86 mm circular bottom puck",GREEN)
card(s,8.8,1.75,3.8,2.0,"User experience","Every input remappable\nBLE and wired USB\n13 individually addressed RGB keys\nClear caps + one ceramic accent",AMBER)
bullets(s,["The supplied visual sketch is an aesthetic reference; the electrical count remains the previously specified 4×4 set of sixteen logical inputs."],1.0,4.45,11.2,1.1,18); footer(s,2)

# 3 — electronics
s = prs.slides.add_slide(prs.slide_layouts[6]); title(s,"Electronics selection","Exact primary components for each of the two devices")
items=[("Compute + radio","Raytac MDBT50Q-1MV2\nnRF52840 certified module"),("Power path","TI BQ24074 charger\nTPS63031 3.3 V buck-boost"),("RGB power","TPS61023 5 V boost\nTPS22918 switch + AHCT buffer"),("Sensors","Alps RKJX21224001 joystick\nAT42QT1010 touch\nMAX17048 fuel gauge"),("Connectivity","GCT USB4105 USB-C\nTI USB ESD + Nexperia VBUS TVS"),("Battery","TinyCircuits ASR00012\nProtected 1,000 mAh; 1 A max")]
for i,(h,b) in enumerate(items): card(s,.7+(i%3)*4.05,1.65+(i//3)*2.25,3.75,1.85,h,b,[BLUE,GREEN,AMBER][i%3])
footer(s,3)

# 4 — quantities
s = prs.slides.add_slide(prs.slide_layouts[6]); title(s,"Exactly two devices","No third assembled unit and no discretionary spare quantities")
card(s,.8,1.7,3.7,2.0,"26 each","Kailh Polia switches\nKailh hot-swap sockets\n1N4148W matrix diodes\n12 mA SK6812MINI-E LEDs",BLUE)
card(s,4.8,1.7,3.7,2.0,"2 each","Radio modules • batteries\nUSB-C • joystick • encoder\nEvery power and sensor IC\nJST-SH battery connectors",GREEN)
card(s,8.8,1.7,3.7,2.0,"Pack minimums","2 × Adafruit 5068 12-cap packs\n1 × Cerakey four-cap pack\n1 silicone sheet + adhesive\n5 bare PCBs may be fab minimum",AMBER)
bullets(s,["57 frozen BOM lines include per-board reference designators, fabricated parts, and reusable programming tools.","JLC attrition quantities are manufacturing consumables—not parts for a third device.","The machine-readable CSV is the purchasing authority for quantities."],1.0,4.2,11.1,1.9,16); footer(s,4)

# 5 — assembly
s = prs.slides.add_slide(prs.slide_layouts[6]); title(s,"Who assembles what","Designed around no at-home reflow")
card(s,.8,1.65,5.65,3.85,"JLCPCB / contract manufacturer","All SMT passives and ICs\nRaytac radio module\nBQ24074 QFN charger\nPower converters and protection\n13 reverse-mounted RGB LEDs\n13 hot-swap sockets + diodes\nUSB-C hybrid assembly preferred",BLUE)
card(s,6.85,1.65,5.65,3.85,"Hand assembly","Install encoder and joystick\nInsert MX switches and keycaps\nConnect protected battery\nInstall heat-set inserts and screws\nFit printed enclosure parts\nCut/apply circular silicone pad\nProgram and run acceptance test",GREEN)
footer(s,5)

# 6 — power/risk
s = prs.slides.add_slide(prs.slide_layouts[6]); title(s,"Power budget and non-negotiable limits")
card(s,.75,1.6,3.85,2.15,"RGB budget","13 × 12 mA channels\n≈169 mA total at 5 V incl. static\n≈313 mA battery draw at 3.0 V\nFirmware cap: 156 mA channels",BLUE)
card(s,4.75,1.6,3.85,2.15,"Battery + charging","1,000 mAh protected cell\n1 A maximum discharge\nTarget total draw <0.45 A\nCharge limited to ≈0.4 A",GREEN)
card(s,8.75,1.6,3.85,2.15,"Thermal constraint","BQ24074 may dissipate ≈0.8 W\nThermal copper/vias required\nClosed-case charge test mandatory\nNo USB-PD claim",AMBER)
bullets(s,["Do not substitute generic 60 mA-per-pixel LEDs.","The RGB rail defaults off and is enabled only after controlled startup."],1.0,4.45,11.2,1.3,18); footer(s,6)

# 7 — cost
s = prs.slides.add_slide(prs.slide_layouts[6]); title(s,"Two-device prototype budget","Planning allowance; final price comes from live JLC and distributor carts")
card(s,.8,1.65,3.7,2.2,"PCBs + assembly","$110–190\nFive bare boards may be the fabrication minimum; populate only two. Includes double-sided setup and ordinary SMT labor.",BLUE)
card(s,4.8,1.65,3.7,2.2,"Controls + mechanics","$90–170\nSwitches, caps, joysticks, encoders, batteries, hardware, printed parts, and foot material.",GREEN)
card(s,8.8,1.65,3.7,2.2,"Shipping + contingency","$50–100\nMultiple vendors, taxes, rare-part loading, and first-article rework allowance.",AMBER)
big=s.shapes.add_textbox(Inches(2.0),Inches(4.45),Inches(9.3),Inches(1.0)); p=big.text_frame.paragraphs[0]; p.text="$250–460 total  •  $125–230 per working prototype"; p.alignment=PP_ALIGN.CENTER; p.font.name="Avenir Next"; p.font.size=Pt(26); p.font.bold=True; p.font.color.rgb=INK
footer(s,7)

# 8 — release gate
s = prs.slides.add_slide(prs.slide_layouts[6]); title(s,"Before placing the complete order","The BOM is selected; product validation still depends on the PCB and physical prototype")
bullets(s,["Complete schematic and four-layer routing; pass ERC/DRC and independent review.","Replace every placeholder with its exact manufacturer footprint and STEP model.","Redesign battery pocket, joystick opening/tail path, USB opening, and insert bosses.","Verify radio antenna keep-out against copper, battery, screws, and circular base.","Validate every JLC MPN, side, rotation, and package in the live order preview.","Build two units; test power, charging temperature, USB, BLE range, RGB, touch, all controls, and enclosure fit."],1.0,1.55,11.4,4.9,18)
footer(s,8)

# 9 — decision
s = prs.slides.add_slide(prs.slide_layouts[6]); title(s,"Recommendation")
card(s,.9,1.65,5.55,3.7,"Buy now","Two batteries for fit checks\nTwo joysticks and encoders\nTwo keycap packs + ceramic pack\nSwitches and mechanical samples\nProgramming cable / nRF52840-DK if needed",GREEN)
card(s,6.85,1.65,5.55,3.7,"Hold until layout validation","JLC assembled PCB order\nFinal printed enclosure batch\nCustom-cut production foot pads\nAny volume purchase beyond the two devices",AMBER)
bullets(s,["Procurement source: hardware/bom.csv", "Engineering rationale and gates: hardware/BOM.md"],1.1,5.75,11.0,.9,15,MUTED); footer(s,9)

OUT.parent.mkdir(parents=True, exist_ok=True)
prs.save(OUT)
print(OUT)
